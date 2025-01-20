import os
import json
import csv
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional, Union, Dict, Any
from dataclasses import dataclass
import traceback
import sys
import logging

try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    from PyPDF2 import PdfReader
    from PIL import Image
    import pytesseract
    from bs4 import BeautifulSoup
except ImportError as e:
    print(f"インポートエラー: {str(e)}")
    print("pip install -r requirements.txtを実行してください。")
    raise

# ログファイルの設定
LOG_FILE = 'markitdown.log'  # ログファイルをカレントディレクトリに保存
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(LOG_FILE, encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class ConversionResult:
    text_content: str
    metadata: Dict[str, Any]

class MarkItDown:
    def __init__(self, llm_client=None, llm_model: str = "gpt-4", debug: bool = False):
        self.llm_client = llm_client
        self.llm_model = llm_model
        self.debug = debug
        self.supported_extensions = {
            '.docx': self._convert_docx,
            '.xlsx': self._convert_xlsx,
            '.pptx': self._convert_pptx,
            '.pdf': self._convert_pdf,
            '.jpg': self._convert_image,
            '.jpeg': self._convert_image,
            '.png': self._convert_image,
            '.gif': self._convert_image,
            '.html': self._convert_html,
            '.csv': self._convert_csv,
            '.json': self._convert_json,
            '.xml': self._convert_xml,
            '.md': self._convert_markdown,
            '.txt': self._convert_text
        }

    def _log(self, message: str):
        """デバッグログを出力"""
        if self.debug:
            logger.debug(message)

    def convert(self, file_path: Union[str, Path]) -> ConversionResult:
        """
        ファイルをMarkdown形式に変換します。
        
        Args:
            file_path: 変換対象ファイルのパス
            
        Returns:
            ConversionResult: 変換結果とメタデータ
        """
        try:
            file_path = Path(file_path).resolve()
            self._log(f"入力ファイル: {file_path}")
            
            if not file_path.exists():
                error_msg = f"ファイルが見つかりません: {file_path}"
                logger.error(error_msg)
                raise FileNotFoundError(error_msg)

            extension = file_path.suffix.lower()
            self._log(f"ファイル拡張子: {extension}")
            
            if extension not in self.supported_extensions:
                error_msg = f"未対応のファイル形式です: {extension}"
                logger.error(error_msg)
                raise ValueError(error_msg)

            converter = self.supported_extensions[extension]
            self._log(f"変換処理開始: {converter.__name__}")
            
            result = converter(file_path)
            self._log("変換処理完了")

            # 出力ファイルの保存
            output_path = file_path.with_suffix('.md')
            if output_path != file_path:
                self.save_markdown(result, output_path)
                print(f"変換完了: {output_path}")

            return result
            
        except Exception as e:
            logger.error(f"エラー発生: {str(e)}")
            logger.exception("詳細なエラー情報:")
            raise

    def save_markdown(self, result: ConversionResult, output_path: Path):
        """
        変換結果をファイルに保存
        """
        self._log(f"保存先: {output_path}")
        try:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(result.text_content)
            self._log("保存完了")
        except Exception as e:
            logger.error(f"保存エラー: {str(e)}")
            raise

    def _convert_docx(self, file_path: Path) -> ConversionResult:
        """Word文書をMarkdownに変換"""
        self._log("Word文書の変換開始")
        doc = Document(file_path)
        markdown = []
        metadata = {"title": doc.core_properties.title or ""}

        for paragraph in doc.paragraphs:
            style = paragraph.style.name
            text = paragraph.text.strip()
            
            if not text:
                continue
                
            if style.startswith('Heading'):
                level = style[-1]
                markdown.append(f"{'#' * int(level)} {text}\n")
            else:
                markdown.append(f"{text}\n\n")

        self._log("Word文書の変換完了")
        return ConversionResult(
            text_content=''.join(markdown),
            metadata=metadata
        )

    def _convert_xlsx(self, file_path: Path) -> ConversionResult:
        """Excelファイルをマークダウンテーブルに変換"""
        self._log("Excelファイルの変換開始")
        wb = load_workbook(file_path)
        markdown = []
        metadata = {"sheets": []}

        for sheet in wb.worksheets:
            markdown.append(f"## {sheet.title}\n\n")
            metadata["sheets"].append(sheet.title)
            
            # ヘッダー行の処理
            headers = []
            for cell in sheet[1]:
                headers.append(str(cell.value or ''))
            
            if headers:
                markdown.append('| ' + ' | '.join(headers) + ' |\n')
                markdown.append('|' + '---|' * len(headers) + '\n')
                
                # データ行の処理
                for row in sheet.iter_rows(min_row=2):
                    row_data = [str(cell.value or '') for cell in row]
                    markdown.append('| ' + ' | '.join(row_data) + ' |\n')
                
                markdown.append('\n')

        self._log("Excelファイルの変換完了")
        return ConversionResult(
            text_content=''.join(markdown),
            metadata=metadata
        )

    def _convert_pptx(self, file_path: Path) -> ConversionResult:
        """PowerPointをMarkdownに変換"""
        self._log("PowerPointの変換開始")
        prs = Presentation(file_path)
        markdown = []
        metadata = {"slide_count": len(prs.slides)}

        for i, slide in enumerate(prs.slides, 1):
            markdown.append(f"# スライド {i}\n\n")
            
            if slide.shapes.title:
                markdown.append(f"## {slide.shapes.title.text}\n\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown.append(f"{shape.text.strip()}\n\n")

        self._log("PowerPointの変換完了")
        return ConversionResult(
            text_content=''.join(markdown),
            metadata=metadata
        )

    def _convert_pdf(self, file_path: Path) -> ConversionResult:
        """PDFをMarkdownに変換"""
        self._log("PDFの変換開始")
        try:
            reader = PdfReader(str(file_path))
            markdown = []
            metadata = {"page_count": len(reader.pages)}

            for i, page in enumerate(reader.pages, 1):
                self._log(f"ページ {i} の処理中")
                markdown.append(f"# ページ {i}\n\n")
                text = page.extract_text()
                if text:
                    # テキストの正規化
                    text = text.replace('\r', '\n')
                    text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
                    markdown.append(f"{text}\n\n")
                else:
                    self._log(f"ページ {i} にテキストが見つかりませんでした")

            self._log("PDFの変換完了")
            return ConversionResult(
                text_content=''.join(markdown),
                metadata=metadata
            )
        except Exception as e:
            logger.error(f"PDF処理エラー: {str(e)}")
            logger.exception("詳細なエラー情報:")
            raise

    def _convert_image(self, file_path: Path) -> ConversionResult:
        """画像をMarkdownに変換（OCRとメタデータ抽出）"""
        self._log("画像の変換開始")
        image = Image.open(file_path)
        metadata = {
            "format": image.format,
            "size": image.size,
            "mode": image.mode
        }

        # EXIFメタデータの抽出
        if hasattr(image, '_getexif'):
            exif = image._getexif()
            if exif:
                metadata["exif"] = {
                    k: str(v) for k, v in exif.items()
                }

        # OCRによるテキスト抽出
        try:
            text = pytesseract.image_to_string(image, lang='jpn+eng')
            markdown = [
                f"# 画像情報\n\n",
                f"![{file_path.name}]({file_path})\n\n",
                "## 抽出テキスト\n\n",
                f"{text}\n\n"
            ]
        except Exception as e:
            self._log(f"OCRエラー: {str(e)}")
            markdown = [
                f"# 画像情報\n\n",
                f"![{file_path.name}]({file_path})\n\n",
                f"OCRエラー: {str(e)}\n\n"
            ]

        self._log("画像の変換完了")
        return ConversionResult(
            text_content=''.join(markdown),
            metadata=metadata
        )

    def _convert_html(self, file_path: Path) -> ConversionResult:
        """HTMLをMarkdownに変換"""
        self._log("HTMLの変換開始")
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            
        markdown = []
        metadata = {
            "title": soup.title.string if soup.title else "",
            "encoding": soup.original_encoding
        }

        # ヘッダーの処理
        for i in range(1, 7):
            for header in soup.find_all(f'h{i}'):
                markdown.append(f"{'#' * i} {header.get_text().strip()}\n\n")

        # 段落の処理
        for p in soup.find_all('p'):
            markdown.append(f"{p.get_text().strip()}\n\n")

        # リストの処理
        for ul in soup.find_all('ul'):
            for li in ul.find_all('li'):
                markdown.append(f"* {li.get_text().strip()}\n")
            markdown.append('\n')

        self._log("HTMLの変換完了")
        return ConversionResult(
            text_content=''.join(markdown),
            metadata=metadata
        )

    def _convert_csv(self, file_path: Path) -> ConversionResult:
        """CSVをMarkdownテーブルに変換"""
        self._log("CSVの変換開始")
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            rows = list(reader)

        markdown = []
        metadata = {"row_count": len(rows)}

        if rows:
            # ヘッダー行
            markdown.append('| ' + ' | '.join(rows[0]) + ' |\n')
            markdown.append('|' + '---|' * len(rows[0]) + '\n')
            
            # データ行
            for row in rows[1:]:
                markdown.append('| ' + ' | '.join(row) + ' |\n')

        self._log("CSVの変換完了")
        return ConversionResult(
            text_content=''.join(markdown),
            metadata=metadata
        )

    def _convert_json(self, file_path: Path) -> ConversionResult:
        """JSONをMarkdownに変換"""
        self._log("JSONの変換開始")
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        def format_json(obj, level=0):
            if isinstance(obj, dict):
                lines = []
                for k, v in obj.items():
                    lines.append(f"{'#' * (level + 1)} {k}\n\n")
                    lines.append(format_json(v, level + 1))
                return ''.join(lines)
            elif isinstance(obj, list):
                return ''.join(f"* {format_json(item, level)}\n" for item in obj)
            else:
                return f"{str(obj)}\n\n"

        self._log("JSONの変換完了")
        return ConversionResult(
            text_content=format_json(data),
            metadata={"type": "json"}
        )

    def _convert_xml(self, file_path: Path) -> ConversionResult:
        """XMLをMarkdownに変換"""
        self._log("XMLの変換開始")
        tree = ET.parse(file_path)
        root = tree.getroot()

        def format_xml(element, level=0):
            lines = []
            lines.append(f"{'#' * (level + 1)} {element.tag}\n\n")
            
            if element.text and element.text.strip():
                lines.append(f"{element.text.strip()}\n\n")
                
            for child in element:
                lines.append(format_xml(child, level + 1))
                
            return ''.join(lines)

        self._log("XMLの変換完了")
        return ConversionResult(
            text_content=format_xml(root),
            metadata={"root_tag": root.tag}
        )

    def _convert_markdown(self, file_path: Path) -> ConversionResult:
        """Markdownファイルをそのまま返す"""
        self._log("Markdownファイルの読み込み")
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return ConversionResult(
            text_content=content,
            metadata={"type": "markdown"}
        )

    def _convert_text(self, file_path: Path) -> ConversionResult:
        """テキストファイルを基本的なMarkdownに変換"""
        self._log("テキストファイルの変換開始")
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # 単純な変換: 空行で段落を分割
        paragraphs = [p.strip() for p in content.split('\n\n')]
        markdown = '\n\n'.join(paragraphs)
        
        return ConversionResult(
            text_content=markdown,
            metadata={"type": "text"}
        )

def main():
    import argparse
    parser = argparse.ArgumentParser(description='ファイルをMarkdown形式に変換')
    parser.add_argument('input_file', help='変換対象ファイル')
    parser.add_argument('-d', '--debug', action='store_true', help='デバッグモードを有効化')
    args = parser.parse_args()

    try:
        converter = MarkItDown(debug=args.debug)
        converter.convert(args.input_file)
    except Exception as e:
        print(f"エラー: {str(e)}", file=sys.stderr)
        if args.debug:
            traceback.print_exc()
        sys.exit(1)

if __name__ == '__main__':
    main()